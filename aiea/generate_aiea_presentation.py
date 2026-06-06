from __future__ import annotations

from io import BytesIO
from pathlib import Path
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
AIEA_DIR = ROOT / "aiea"
REFERENCE_PPT = ROOT / "thesis" / "答辩" / "论文汇报PPT_模板风格版.pptx"
OUTPUT_PPT = AIEA_DIR / "aiea_conference_presentation.pptx"

SLIDE_WIDTH = 12179300
SLIDE_HEIGHT = 6858000

FONT_CN = "Noto Sans CJK SC"
FONT_EN = "Calibri"

RED = RGBColor(155, 31, 35)
RED_DARK = RGBColor(122, 20, 23)
TEXT = RGBColor(48, 48, 48)
TEXT_LIGHT = RGBColor(92, 92, 92)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(31, 96, 156)
TEAL = RGBColor(0, 138, 150)
GREEN = RGBColor(28, 145, 93)
ORANGE = RGBColor(220, 116, 26)
LIGHT_BLUE = RGBColor(231, 239, 248)
LIGHT_GRAY = RGBColor(247, 247, 247)
MID_GRAY = RGBColor(128, 128, 128)

PRESENTATION_NAME = "Choice-Head Distillation for Efficient Dental MCQ Answering"
FOOTER_NAME = "AIEA 2026 Conference Presentation"

SPEAKER_NOTES = [
    [
        "各位老师、各位专家好。我今天汇报的题目是用于高效牙科选择题自动答题的 Choice-Head 蒸馏方法。",
        "这个工作的核心目标，是把大型医疗语言模型的答题能力迁移到更小、更容易部署的学生模型上。",
    ],
    [
        "这一页先给各位一个整体路线图。",
        "我会先讲问题背景和研究问题，再介绍 Choice-Head 方法，然后给出实验设置、核心结果、两点贡献，以及最后的讨论和结论。",
    ],
    [
        "医疗大语言模型在考试类基准上已经很强，但部署成本依然很高，包括推理成本、硬件需求和复现成本。",
        "而我们的任务不是开放式聊天，而是标准的五选一牙科选择题，所以蒸馏目标未必还要沿用全词表那套通用做法。",
    ],
    [
        "这里最核心的问题不是怎样完整模仿一个大语言模型，而是对于五选一考试题，到底什么信息最值得迁移。",
        "学生究竟要学整个词表分布，还是只学最后真正影响答题的决策结构，同时还要兼容 API 黑盒教师。",
    ],
    [
        "我们的方法叫 Choice-Head distillation。",
        "核心思想是只蒸馏 A、B、C、D、E 五个选项上的概率分布，而不是让学生去拟合教师的完整词表输出。",
        "在第一阶段里，我们用选项分布上的 KL 散度加标准答案上的交叉熵共同训练学生。",
    ],
    [
        "这个设计有三个关键价值。第一，它让训练目标和任务本身对齐。第二，它去掉了大量无关监督。第三，它让黑盒强教师也能纳入蒸馏流程。",
        "换句话说，我们只保留真正影响最后选项决策的教师信号。",
    ],
    [
        "实验使用的是基于 CMExam 构建的全量重分割数据，一共 6591 道单选医学题。",
        "主教师是 DeepSeek-V3，学生是 Qwen2.5-7B 和 Qwen2.5-14B。表现最好的 14B 配置只使用一阶段 Choice-Head 蒸馏，并采用 LoRA 微调。",
    ],
    [
        "核心结果非常直接。14B 学生的零样本基线是 83.55%，教师是 87.18%。",
        "经过 Choice-Head 蒸馏后，14B 学生的平均准确率达到 88.67%，最佳达到 89.10%。",
        "这里的 Full Test Accuracy 指的是 991 题全量测试集，Dental Test Accuracy 指的是 125 题牙科子集。",
        "Zero-shot baseline 就是不蒸馏直接测试，Stage 1 mean 是多次蒸馏运行的平均值，Stage 1 best 是其中最好的单次结果。",
        "也就是说，学生不仅明显超过了自己的零样本基线，而且在同一个 991 题测试集上也超过了教师。",
    ],
    [
        "第一个贡献是方法上的：我们说明了对于结构化医学选择题，蒸馏目标应该围绕任务决策结构来定义。",
        "第二个贡献是实证上的：在更大、更可靠的测试集上，一个更小的学生模型在这种任务对齐的设置下可以超过更强的教师模型。",
    ],
    [
        "一个很重要的实际发现是，更复杂的训练流程不一定更好。",
        "在我们的实验里，最强的 14B 结果恰恰来自只做 Stage 1。",
        "更大的结论是，蒸馏应该围绕任务结构来设计，而不是机械沿用通用大模型训练范式。",
    ],
    [
        "最后，这项工作说明了一件事情：想实现高效部署，并不一定要重现教师模型的全部语言行为。",
        "对于结构化医学选择题任务，只要迁移选项级决策分布，就可能获得很强的效果。",
        "谢谢各位老师和专家聆听，欢迎批评指正。",
        "Q&A 可补充三点：为什么只蒸馏五个选项、为什么学生能超过教师、以及这个方法为什么不能直接照搬到开放式问答。",
    ],
    [
        "这一页是我为现场问答准备的备用页。",
        "如果被问到为什么只蒸馏五个选项，我会强调下游任务本质上就是五选一决策问题。",
        "如果被问到为什么学生能超过教师，我会强调学生同时学习了教师软标签结构和任务本身的优化目标。",
        "如果被问到能否扩展到开放式问答，我会说明当前方法不能直接照搬，因为开放式任务需要重新定义监督目标。",
    ],
]


def set_text_style(paragraph, font_size, color=TEXT, bold=False, font_name=FONT_EN):
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=20,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name=FONT_EN,
    vertical_anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = vertical_anchor
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    set_text_style(paragraph, font_size, color=color, bold=bold, font_name=font_name)
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items,
    *,
    font_size=20,
    color=TEXT_LIGHT,
    font_name=FONT_EN,
    level0_indent=0.28,
    space_after=8,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.clear()
    for index, item in enumerate(items):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.left_margin = Inches(level0_indent)
        paragraph.space_after = Pt(space_after)
        paragraph.alignment = PP_ALIGN.LEFT
        set_text_style(paragraph, font_size, color=color, font_name=font_name)
    return box


def add_footer(slide, index, total, title=FOOTER_NAME):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        SLIDE_HEIGHT - Inches(0.5),
        SLIDE_WIDTH,
        Inches(0.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()

    add_textbox(
        slide,
        Inches(4.05),
        SLIDE_HEIGHT - Inches(0.39),
        Inches(4.4),
        Inches(0.25),
        title,
        font_size=15,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        SLIDE_WIDTH - Inches(0.78),
        SLIDE_HEIGHT - Inches(0.39),
        Inches(0.5),
        Inches(0.25),
        f"{index}/{total}",
        font_size=15,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
        font_name=FONT_EN,
    )


def add_title(slide, number, title, *, title_color=TEXT):
    add_textbox(
        slide,
        Inches(0.72),
        Inches(0.38),
        Inches(0.8),
        Inches(0.5),
        number,
        font_size=30,
        color=RED,
        font_name=FONT_EN,
    )
    add_textbox(
        slide,
        Inches(1.7),
        Inches(0.34),
        Inches(8.8),
        Inches(0.6),
        title,
        font_size=28,
        color=title_color,
        bold=True,
    )
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.82),
        Inches(1.14),
        Inches(2.05),
        Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()


def add_round_box(slide, left, top, width, height, *, fill, line):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.4)
    return shape


def add_arrow_block(slide, left, top, width, height, text, *, fill, color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    set_text_style(paragraph, 15, color=color, bold=True)
    return shape


def set_slide_notes(slide, paragraphs):
    notes_slide = slide.notes_slide
    for shape in notes_slide.shapes:
        if shape.is_placeholder and hasattr(shape, "text_frame") and shape.placeholder_format.idx == 2:
            text_frame = shape.text_frame
            text_frame.clear()
            for index, text in enumerate(paragraphs):
                paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                paragraph.text = text
                set_text_style(paragraph, 12, color=TEXT, font_name=FONT_CN)
            return


def add_card(slide, left, top, width, height, title, body, *, line_color, title_color=None):
    card = add_round_box(slide, left, top, width, height, fill=WHITE, line=line_color)
    text_frame = card.text_frame
    text_frame.clear()
    title_paragraph = text_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.space_after = Pt(6)
    set_text_style(title_paragraph, 18, color=title_color or line_color, bold=True)

    for line in body:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line
        paragraph.space_after = Pt(4)
        set_text_style(paragraph, 13, color=TEXT_LIGHT)
    return card


def add_metric_card(slide, left, top, width, height, title, value, note, *, accent):
    card = add_round_box(slide, left, top, width, height, fill=LIGHT_GRAY, line=accent)
    text_frame = card.text_frame
    text_frame.clear()

    title_paragraph = text_frame.paragraphs[0]
    title_paragraph.text = title
    set_text_style(title_paragraph, 14, color=accent, bold=True)

    value_paragraph = text_frame.add_paragraph()
    value_paragraph.text = value
    value_paragraph.space_before = Pt(5)
    set_text_style(value_paragraph, 24, color=TEXT, bold=True, font_name=FONT_EN)

    note_paragraph = text_frame.add_paragraph()
    note_paragraph.text = note
    note_paragraph.space_before = Pt(2)
    set_text_style(note_paragraph, 11, color=TEXT_LIGHT)
    return card


def add_logo(slide):
    if not REFERENCE_PPT.exists():
        return
    try:
        with zipfile.ZipFile(REFERENCE_PPT) as archive:
            media_names = [name for name in archive.namelist() if name.startswith("ppt/media/")]
            if not media_names:
                return
            data = archive.read(media_names[0])
        slide.shapes.add_picture(BytesIO(data), Inches(0.35), Inches(0.16), width=Inches(2.1))
    except (KeyError, zipfile.BadZipFile):
        return


def add_bar_comparison(slide, left, top, width, height, items):
    chart_box = add_round_box(slide, left, top, width, height, fill=WHITE, line=BLUE)
    chart_box.text_frame.clear()

    max_value = max(item[1] for item in items)
    min_value = min(item[1] for item in items)
    value_span = max_value - min_value
    if value_span <= 0:
        value_span = 1.0

    inner_left = left + Inches(0.35)
    label_width = Inches(1.8)
    bar_left = inner_left + label_width
    bar_top = top + Inches(0.38)
    bar_max_width = width - Inches(2.15)
    bar_height = Inches(0.33)
    row_gap = Inches(0.54)

    for index, (label, value, color) in enumerate(items):
        current_top = bar_top + index * row_gap
        add_textbox(
            slide,
            inner_left,
            current_top - Inches(0.02),
            label_width - Inches(0.15),
            Inches(0.25),
            label,
            font_size=13,
            color=TEXT,
        )

        baseline = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            bar_left,
            current_top,
            bar_max_width,
            bar_height,
        )
        baseline.fill.solid()
        baseline.fill.fore_color.rgb = LIGHT_BLUE
        baseline.line.fill.background()

        normalized = 0.42 + 0.58 * ((value - min_value) / value_span)
        value_width = int(bar_max_width * normalized)
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            bar_left,
            current_top,
            value_width,
            bar_height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        add_textbox(
            slide,
            bar_left + bar_max_width + Inches(0.08),
            current_top - Inches(0.02),
            Inches(0.75),
            Inches(0.25),
            f"{value:.2f}%",
            font_size=13,
            color=color,
            bold=True,
            font_name=FONT_EN,
        )


def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(1.35),
        Inches(10.5),
        Inches(1.1),
        PRESENTATION_NAME,
        font_size=28,
        color=TEXT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(2.0),
        Inches(2.5),
        Inches(8.6),
        Inches(0.42),
        "Decision-Space Distillation for Deployable Medical QA",
        font_size=18,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
        font_name=FONT_EN,
    )
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(4.4),
        Inches(3.1),
        Inches(3.25),
        Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    add_textbox(
        slide,
        Inches(2.0),
        Inches(3.48),
        Inches(8.9),
        Inches(0.5),
        "DeepSeek-V3 Teacher  |  Qwen2.5-14B Student  |  CMExam Full-Test 991 Questions",
        font_size=16,
        color=TEXT_LIGHT,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(2.55),
        Inches(4.02),
        Inches(7.8),
        Inches(0.5),
        "Key Result: 89.10% student accuracy > 87.18% teacher accuracy",
        font_size=19,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name=FONT_EN,
    )

    info = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        SLIDE_HEIGHT - Inches(1.95),
        SLIDE_WIDTH,
        Inches(1.8),
    )
    info.fill.solid()
    info.fill.fore_color.rgb = RED
    info.line.fill.background()
    add_textbox(slide, Inches(1.9), Inches(5.72), Inches(3.0), Inches(0.35), "Presenter: Tianyuan Chen", font_size=21, color=WHITE, bold=True)
    add_textbox(slide, Inches(7.15), Inches(5.72), Inches(3.4), Inches(0.35), "Supervisor: Tichao Xiong", font_size=21, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.9), Inches(6.36), Inches(3.3), Inches(0.35), "Program: Applied Artificial Intelligence", font_size=19, color=WHITE, bold=True)
    add_textbox(slide, Inches(7.15), Inches(6.36), Inches(2.9), Inches(0.35), "AIEA 2026, Shenzhen", font_size=19, color=WHITE, bold=True)

    bottom = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, SLIDE_HEIGHT - Inches(0.48), SLIDE_WIDTH, Inches(0.43))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = RED_DARK
    bottom.line.fill.background()
    add_textbox(
        slide,
        Inches(5.25),
        Inches(7.07),
        Inches(2.8),
        Inches(0.2),
        "June 2026",
        font_size=15,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        font_name=FONT_EN,
    )


def build_outline(slide, total):
    add_textbox(slide, Inches(0.95), Inches(0.62), Inches(2.6), Inches(0.55), "Presentation Outline", font_size=30, color=TEXT, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.86), Inches(1.2), Inches(1.95), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()

    outline_box = add_round_box(slide, Inches(0.95), Inches(1.7), Inches(10.15), Inches(4.65), fill=WHITE, line=TEAL)
    outline_box.text_frame.clear()
    add_bullets(
        slide,
        Inches(1.25),
        Inches(1.98),
        Inches(9.4),
        Inches(4.1),
        [
            "1. Motivation and problem setting",
            "2. Research questions and Choice-Head method",
            "3. Why decision-space distillation matters",
            "4. Experimental setup and benchmark scale",
            "5. Main results and student-over-teacher finding",
            "6. Contributions, discussion, and conclusion",
        ],
        font_size=22,
        color=TEXT,
        level0_indent=0.16,
        space_after=10,
    )
    add_footer(slide, 2, total)


def build_motivation(slide, total):
    add_title(slide, "01", "Motivation and Problem Setting", title_color=RED)
    banner = add_round_box(slide, Inches(0.85), Inches(1.6), Inches(11.1), Inches(1.0), fill=LIGHT_BLUE, line=TEAL)
    banner_tf = banner.text_frame
    banner_tf.clear()
    paragraph = banner_tf.paragraphs[0]
    paragraph.text = "Goal: preserve exam performance while making dental MCQ systems smaller, cheaper, and easier to deploy."
    paragraph.alignment = PP_ALIGN.CENTER
    set_text_style(paragraph, 21, color=BLUE, bold=True)

    add_card(
        slide,
        Inches(0.9),
        Inches(3.0),
        Inches(3.45),
        Inches(2.45),
        "Current Tension",
        [
            "Medical LLMs are accurate on exam benchmarks",
            "But inference cost and hardware demand remain high",
            "Deployment value matters beyond raw benchmark scores",
        ],
        line_color=RED,
    )
    add_card(
        slide,
        Inches(4.65),
        Inches(3.0),
        Inches(3.15),
        Inches(2.45),
        "Task Structure",
        [
            "Dental QA is a five-option MCQ task",
            "Decision space is small and well-structured",
            "Open-ended generation is not required here",
        ],
        line_color=TEAL,
    )
    add_card(
        slide,
        Inches(8.1),
        Inches(3.0),
        Inches(3.0),
        Inches(2.45),
        "Mismatch",
        [
            "Full-vocabulary distillation is redundant",
            "API teachers rarely expose internal logits",
            "Task-aligned supervision is needed",
        ],
        line_color=ORANGE,
    )
    add_footer(slide, 2, total)


def build_research_question(slide, total):
    add_title(slide, "02", "Research Questions")
    add_metric_card(slide, Inches(0.85), Inches(1.6), Inches(3.1), Inches(1.2), "RQ1", "What to Distill?", "Whole vocabulary or only A/B/C/D/E?", accent=TEAL)
    add_metric_card(slide, Inches(4.25), Inches(1.6), Inches(3.1), Inches(1.2), "RQ2", "Can Students Exceed Teachers?", "Under a task-aligned objective", accent=GREEN)
    add_metric_card(slide, Inches(7.65), Inches(1.6), Inches(3.1), Inches(1.2), "RQ3", "Can It Stay Black-Box Friendly?", "Compatible with API teachers", accent=ORANGE)

    box = add_round_box(slide, Inches(0.95), Inches(3.25), Inches(10.4), Inches(2.4), fill=WHITE, line=BLUE)
    text_frame = box.text_frame
    text_frame.clear()
    lead = text_frame.paragraphs[0]
    lead.text = "Central research question"
    set_text_style(lead, 22, color=BLUE, bold=True)
    for line in [
        "For a five-option medical exam task, the important transfer target may be the decision structure rather than the full language-model distribution.",
        "If the supervision target is redesigned around the downstream task, smaller students may become both deployable and unexpectedly strong.",
    ]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line
        paragraph.space_after = Pt(6)
        set_text_style(paragraph, 16, color=TEXT_LIGHT)
    add_footer(slide, 3, total)


def build_method(slide, total):
    add_title(slide, "03", "Core Method: Choice-Head Distillation")
    add_card(slide, Inches(0.8), Inches(1.5), Inches(2.35), Inches(1.55), "Teacher", ["DeepSeek-V3", "Can be local or API-based", "Provides option-level preference"], line_color=RED)
    add_arrow_block(slide, Inches(3.35), Inches(1.82), Inches(1.25), Inches(0.62), "A/B/C/D/E", fill=TEAL)
    add_card(slide, Inches(4.9), Inches(1.5), Inches(2.2), Inches(1.55), "Choice Head", ["Five-option distribution", "Decision-space supervision", "No full vocabulary needed"], line_color=TEAL)
    add_arrow_block(slide, Inches(7.28), Inches(1.82), Inches(1.15), Inches(0.62), "KL + CE", fill=ORANGE)
    add_card(slide, Inches(8.7), Inches(1.5), Inches(2.45), Inches(1.55), "Student", ["Qwen2.5-7B / 14B", "Learns the final choice boundary", "More deployable"], line_color=GREEN)

    left = add_round_box(slide, Inches(0.85), Inches(3.45), Inches(6.15), Inches(2.25), fill=LIGHT_GRAY, line=TEAL)
    left_tf = left.text_frame
    left_tf.clear()
    title_paragraph = left_tf.paragraphs[0]
    title_paragraph.text = "Stage 1 objective"
    set_text_style(title_paragraph, 22, color=TEAL, bold=True)
    formula_paragraph = left_tf.add_paragraph()
    formula_paragraph.text = "L = alpha * KL(p_T || p_S) + (1 - alpha) * CE"
    formula_paragraph.space_before = Pt(8)
    set_text_style(formula_paragraph, 24, color=TEXT, bold=True, font_name=FONT_EN)
    note_paragraph = left_tf.add_paragraph()
    note_paragraph.text = "The student learns relative preferences across the five answer options instead of the entire token vocabulary."
    note_paragraph.space_before = Pt(6)
    set_text_style(note_paragraph, 15, color=TEXT_LIGHT)
    second_note = left_tf.add_paragraph()
    second_note.text = "Main setting: one Stage 1 epoch with alpha = 0.35 for the best 14B run."
    second_note.space_before = Pt(3)
    set_text_style(second_note, 14, color=TEXT_LIGHT)

    right = add_round_box(slide, Inches(7.25), Inches(3.45), Inches(4.0), Inches(2.25), fill=LIGHT_BLUE, line=BLUE)
    right_tf = right.text_frame
    right_tf.clear()
    title_paragraph = right_tf.paragraphs[0]
    title_paragraph.text = "Design intuition"
    set_text_style(title_paragraph, 21, color=BLUE, bold=True)
    for line in [
        "Keep the part of teacher uncertainty that directly affects the final option.",
        "Drop supervision dimensions that do not matter for a five-choice decision.",
    ]:
        paragraph = right_tf.add_paragraph()
        paragraph.text = line
        paragraph.bullet = True
        paragraph.left_margin = Inches(0.23)
        set_text_style(paragraph, 15, color=TEXT_LIGHT)
    final_line = right_tf.add_paragraph()
    final_line.text = "This is why the method stays lightweight while remaining competitive."
    final_line.space_before = Pt(5)
    set_text_style(final_line, 14, color=BLUE, bold=True)
    add_footer(slide, 5, total)


def build_why_it_matters(slide, total):
    add_title(slide, "04", "Why Choice-Head Matters")
    add_metric_card(slide, Inches(0.85), Inches(1.6), Inches(2.4), Inches(1.25), "Alignment", "Task-Specific", "Supervision follows the five-option structure", accent=TEAL)
    add_metric_card(slide, Inches(3.55), Inches(1.6), Inches(2.4), Inches(1.25), "Efficiency", "Lower Cost", "Less redundant training and memory usage", accent=ORANGE)
    add_metric_card(slide, Inches(6.25), Inches(1.6), Inches(2.4), Inches(1.25), "Teacher Access", "Black-Box Friendly", "Works with API-only teachers", accent=BLUE)
    add_metric_card(slide, Inches(8.95), Inches(1.6), Inches(2.4), Inches(1.25), "Transfer Signal", "Useful Uncertainty", "Preserves option-level confusion structure", accent=GREEN)

    highlight = add_round_box(slide, Inches(1.05), Inches(3.35), Inches(10.0), Inches(1.1), fill=LIGHT_BLUE, line=TEAL)
    highlight_tf = highlight.text_frame
    highlight_tf.clear()
    paragraph = highlight_tf.paragraphs[0]
    paragraph.text = "The method changes distillation from a generic language-model compression problem into a task-specific decision transfer problem."
    paragraph.alignment = PP_ALIGN.CENTER
    set_text_style(paragraph, 20, color=BLUE, bold=True)

    add_bullets(
        slide,
        Inches(1.2),
        Inches(4.8),
        Inches(9.6),
        Inches(1.4),
        [
            "Better aligned supervision often matters more than a more complicated pipeline.",
            "For structured MCQs, the decision-space target is a cleaner and more deployable abstraction.",
        ],
        font_size=18,
        color=TEXT,
        level0_indent=0.2,
    )
    add_footer(slide, 6, total)


def build_setup(slide, total):
    add_title(slide, "05", "Experimental Setup")
    add_metric_card(slide, Inches(0.85), Inches(1.55), Inches(2.45), Inches(1.2), "Dataset", "6,591 Questions", "CMExam-based single-choice resplit", accent=TEAL)
    add_metric_card(slide, Inches(3.5), Inches(1.55), Inches(2.45), Inches(1.2), "Splits", "4608 / 991 / 991", "Train / Val / Test", accent=ORANGE)
    add_metric_card(slide, Inches(6.15), Inches(1.55), Inches(2.45), Inches(1.2), "Students", "Qwen2.5 7B / 14B", "LoRA fine-tuning", accent=GREEN)
    add_metric_card(slide, Inches(8.8), Inches(1.55), Inches(2.45), Inches(1.2), "Teacher", "DeepSeek-V3", "87.18% on the 991-question test", accent=BLUE)

    left = add_round_box(slide, Inches(0.85), Inches(3.05), Inches(5.2), Inches(2.6), fill=WHITE, line=TEAL)
    left_tf = left.text_frame
    left_tf.clear()
    paragraph = left_tf.paragraphs[0]
    paragraph.text = "Evaluation scope"
    set_text_style(paragraph, 20, color=TEAL, bold=True)
    for line in [
        "Main benchmark: 991-question full-data test set",
        "Dental subset retained for specialty-focused validation",
        "Large-scale resplit makes the main conclusion more stable",
    ]:
        bullet = left_tf.add_paragraph()
        bullet.text = line
        bullet.bullet = True
        bullet.left_margin = Inches(0.24)
        set_text_style(bullet, 16, color=TEXT_LIGHT)

    right = add_round_box(slide, Inches(6.3), Inches(3.05), Inches(5.0), Inches(2.6), fill=WHITE, line=ORANGE)
    right_tf = right.text_frame
    right_tf.clear()
    paragraph = right_tf.paragraphs[0]
    paragraph.text = "Training setup"
    set_text_style(paragraph, 20, color=ORANGE, bold=True)
    for line in [
        "LoRA rank = 16, LoRA alpha = 32",
        "Best 14B setting uses one Stage 1 epoch",
        "Main learning rate: 1e-4",
        "Stage 2 is tested but not required for the best 14B result",
    ]:
        bullet = right_tf.add_paragraph()
        bullet.text = line
        bullet.bullet = True
        bullet.left_margin = Inches(0.24)
        set_text_style(bullet, 16, color=TEXT_LIGHT)
    add_footer(slide, 7, total)


def build_results(slide, total):
    add_title(slide, "06", "Main Results")
    add_metric_card(slide, Inches(0.85), Inches(1.55), Inches(2.45), Inches(1.25), "Teacher", "87.18%", "DeepSeek-V3 on full-data test", accent=RED)
    add_metric_card(slide, Inches(3.55), Inches(1.55), Inches(2.45), Inches(1.25), "14B Zero-Shot", "83.55%", "Base student without distillation", accent=MID_GRAY)
    add_metric_card(slide, Inches(6.25), Inches(1.55), Inches(2.45), Inches(1.25), "14B Distilled Mean", "88.67%", "Three-seed average", accent=BLUE)
    add_metric_card(slide, Inches(8.95), Inches(1.55), Inches(2.45), Inches(1.25), "14B Distilled Best", "89.10%", "Best single run", accent=GREEN)

    add_bar_comparison(
        slide,
        Inches(0.9),
        Inches(3.1),
        Inches(6.2),
        Inches(2.65),
        [
            ("Teacher", 87.18, RED),
            ("14B Zero-shot", 83.55, MID_GRAY),
            ("14B Mean", 88.67, BLUE),
            ("14B Best", 89.10, GREEN),
        ],
    )

    right = add_round_box(slide, Inches(7.35), Inches(3.1), Inches(4.0), Inches(1.7), fill=LIGHT_BLUE, line=GREEN)
    right_tf = right.text_frame
    right_tf.clear()
    paragraph = right_tf.paragraphs[0]
    paragraph.text = "Takeaway"
    set_text_style(paragraph, 21, color=GREEN, bold=True)
    for line in [
        "The distilled student gains 5.12 percentage points over the 14B zero-shot baseline on average.",
        "The best 14B student exceeds the teacher on the same 991-question benchmark.",
        "The 7B student also improves strongly, so the method is not tied to one model size.",
    ]:
        bullet = right_tf.add_paragraph()
        bullet.text = line
        bullet.bullet = True
        bullet.left_margin = Inches(0.22)
        set_text_style(bullet, 13, color=TEXT_LIGHT)

    legend = add_round_box(slide, Inches(7.35), Inches(4.98), Inches(4.0), Inches(0.95), fill=WHITE, line=ORANGE)
    legend_tf = legend.text_frame
    legend_tf.clear()
    legend_title = legend_tf.paragraphs[0]
    legend_title.text = "How to read Table 1"
    set_text_style(legend_title, 15, color=ORANGE, bold=True)
    for line in [
        "Full Test = 991-question main benchmark; Dental Test = 125-question subset.",
        "Zero-shot = no distillation; Mean = average runs; Best = strongest single run.",
    ]:
        legend_p = legend_tf.add_paragraph()
        legend_p.text = line
        set_text_style(legend_p, 11, color=TEXT_LIGHT)

    win_badge = add_round_box(slide, Inches(7.82), Inches(5.98), Inches(3.0), Inches(0.44), fill=GREEN, line=GREEN)
    badge_tf = win_badge.text_frame
    badge_tf.clear()
    badge_p = badge_tf.paragraphs[0]
    badge_p.text = "+1.92 points vs teacher"
    badge_p.alignment = PP_ALIGN.CENTER
    set_text_style(badge_p, 16, color=WHITE, bold=True, font_name=FONT_EN)
    add_footer(slide, 8, total)


def build_contributions(slide, total):
    add_title(slide, "07", "Two Main Contributions")
    left = add_round_box(slide, Inches(0.9), Inches(1.75), Inches(4.95), Inches(3.4), fill=WHITE, line=TEAL)
    left_tf = left.text_frame
    left_tf.clear()
    paragraph = left_tf.paragraphs[0]
    paragraph.text = "Contribution 1: Method"
    set_text_style(paragraph, 23, color=TEAL, bold=True)
    for line in [
        "Introduces a task-aligned decision-space distillation framework for five-option medical MCQs.",
        "Replaces generic vocabulary-space supervision with a cleaner option-level target.",
        "Keeps the framework practical for both local and API teachers.",
    ]:
        bullet = left_tf.add_paragraph()
        bullet.text = line
        bullet.bullet = True
        bullet.left_margin = Inches(0.24)
        set_text_style(bullet, 16, color=TEXT_LIGHT)

    right = add_round_box(slide, Inches(6.1), Inches(1.75), Inches(5.0), Inches(3.4), fill=WHITE, line=GREEN)
    right_tf = right.text_frame
    right_tf.clear()
    paragraph = right_tf.paragraphs[0]
    paragraph.text = "Contribution 2: Empirical Finding"
    set_text_style(paragraph, 23, color=GREEN, bold=True)
    for line in [
        "Shows that a smaller student can outperform a stronger teacher under the task-aligned formulation.",
        "Validates the result on a 991-question test set rather than only a tiny benchmark slice.",
        "Frames the student-over-teacher effect as a consequence of better task alignment, not a distillation failure.",
    ]:
        bullet = right_tf.add_paragraph()
        bullet.text = line
        bullet.bullet = True
        bullet.left_margin = Inches(0.24)
        set_text_style(bullet, 16, color=TEXT_LIGHT)
    add_footer(slide, 9, total)


def build_discussion(slide, total):
    add_title(slide, "08", "Discussion and Interpretation")
    add_card(slide, Inches(0.8), Inches(1.6), Inches(3.35), Inches(2.05), "Stage 2 is Not Always Better", ["Helpful for some smaller students", "Can mildly hurt a strong 14B student", "Extra GT fine-tuning may erase soft-label structure"], line_color=ORANGE)
    add_card(slide, Inches(4.35), Inches(1.6), Inches(3.35), Inches(2.05), "Why the Student Can Win", ["Teacher soft labels provide structure", "Student is still optimized for the benchmark", "Task-aligned supervision sharpens the decision boundary"], line_color=GREEN)
    add_card(slide, Inches(7.9), Inches(1.6), Inches(3.35), Inches(2.05), "Broader Message", ["More complexity is not the point", "The supervision target must fit the task", "Decision-space transfer can be a practical design rule"], line_color=BLUE)

    highlight = add_round_box(slide, Inches(0.95), Inches(4.15), Inches(10.3), Inches(1.45), fill=LIGHT_BLUE, line=TEAL)
    highlight_tf = highlight.text_frame
    highlight_tf.clear()
    paragraph = highlight_tf.paragraphs[0]
    paragraph.text = "For structured medical multiple-choice tasks, the critical design choice is not how much of the teacher to copy, but which part of the teacher signal actually matters."
    paragraph.alignment = PP_ALIGN.CENTER
    set_text_style(paragraph, 18, color=BLUE, bold=True)
    add_footer(slide, 10, total)


def build_conclusion(slide, total):
    add_title(slide, "09", "Conclusion")
    summary = add_round_box(slide, Inches(1.05), Inches(1.65), Inches(10.0), Inches(1.2), fill=LIGHT_BLUE, line=TEAL)
    summary_tf = summary.text_frame
    summary_tf.clear()
    paragraph = summary_tf.paragraphs[0]
    paragraph.text = "Choice-Head distillation makes smaller dental MCQ models both practical and highly competitive."
    paragraph.alignment = PP_ALIGN.CENTER
    set_text_style(paragraph, 22, color=BLUE, bold=True)

    add_bullets(
        slide,
        Inches(1.3),
        Inches(3.25),
        Inches(9.2),
        Inches(2.1),
        [
            "Decision-space supervision is an effective alternative to full-vocabulary distillation for structured MCQs.",
            "The best 14B student reaches 89.10% and surpasses the 87.18% DeepSeek-V3 teacher on the 991-question test set.",
            "The framework stays lightweight, deployment-oriented, and compatible with black-box teachers.",
        ],
        font_size=19,
        color=TEXT,
        level0_indent=0.22,
    )
    add_textbox(
        slide,
        Inches(3.35),
        Inches(5.75),
        Inches(5.2),
        Inches(0.45),
        "Thank you. Questions are welcome.",
        font_size=24,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name=FONT_EN,
    )
    add_footer(slide, 11, total)


def build_qa_backup(slide, total):
    add_title(slide, "10", "Q&A Backup")
    add_card(
        slide,
        Inches(0.85),
        Inches(1.55),
        Inches(3.35),
        Inches(3.7),
        "Why distill only five options?",
        [
            "The downstream task is a five-option decision problem.",
            "Full-vocabulary supervision introduces unnecessary dimensions.",
            "Option-level transfer is more task-aligned and more efficient.",
        ],
        line_color=TEAL,
    )
    add_card(
        slide,
        Inches(4.45),
        Inches(1.55),
        Inches(3.35),
        Inches(3.7),
        "Why can the student outperform the teacher?",
        [
            "The student learns teacher soft-label structure.",
            "It is still optimized directly for the benchmark objective.",
            "Task-aligned supervision can yield a stronger decision boundary.",
        ],
        line_color=GREEN,
    )
    add_card(
        slide,
        Inches(8.05),
        Inches(1.55),
        Inches(3.1),
        Inches(3.7),
        "Can this extend to open-ended QA?",
        [
            "Not directly.",
            "The current method is built for structured MCQs.",
            "Open-ended QA would require a different supervision target.",
        ],
        line_color=ORANGE,
    )
    add_footer(slide, 12, total)


def main():
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    total = 12
    build_cover(presentation)
    set_slide_notes(presentation.slides[0], SPEAKER_NOTES[0])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_outline(slide, total)
    set_slide_notes(slide, SPEAKER_NOTES[1])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_motivation(slide, total)
    set_slide_notes(slide, SPEAKER_NOTES[2])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_research_question(slide, total)
    set_slide_notes(slide, SPEAKER_NOTES[3])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_method(slide, total)
    set_slide_notes(slide, SPEAKER_NOTES[4])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_why_it_matters(slide, total)
    set_slide_notes(slide, SPEAKER_NOTES[5])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_setup(slide, total)
    set_slide_notes(slide, SPEAKER_NOTES[6])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_results(slide, total)
    set_slide_notes(slide, SPEAKER_NOTES[7])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_contributions(slide, total)
    set_slide_notes(slide, SPEAKER_NOTES[8])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_discussion(slide, total)
    set_slide_notes(slide, SPEAKER_NOTES[9])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_conclusion(slide, total)
    set_slide_notes(slide, SPEAKER_NOTES[10])

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build_qa_backup(slide, total)
    set_slide_notes(slide, SPEAKER_NOTES[11])

    AIEA_DIR.mkdir(parents=True, exist_ok=True)
    presentation.save(OUTPUT_PPT)
    print(f"Generated: {OUTPUT_PPT}")


if __name__ == "__main__":
    main()