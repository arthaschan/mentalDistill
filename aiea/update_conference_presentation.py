#!/usr/bin/env python3
"""Update AIEA conference PPT to match approved 1830 paper."""

from __future__ import annotations

import shutil
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path("/Users/arthas/git/mentalDistill/aiea")
SRC = BASE / "aiea_conference_presentation.pptx"
PAPER = BASE / "aiea_DentalMCQ Distillation 2026-06-12 1830 英文版.docx"
FIG_OUT = BASE / "_ppt_figure1.png"
OUT = BASE / "aiea_conference_presentation.pptx"
BACKUP = BASE / "aiea_conference_presentation_backup_2026-06-12.pptx"


def extract_paper_figure() -> Path:
    with ZipFile(PAPER) as zf:
        FIG_OUT.write_bytes(zf.read("word/media/image1.png"))
    return FIG_OUT


def replace_text_in_slide(slide, old: str, new: str) -> int:
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if old in shape.text:
            shape.text = shape.text.replace(old, new)
            count += 1
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    count += 1
    return count


def set_footer(slide, page_num: int) -> None:
    target = f"{page_num}/12"
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip().endswith("/12"):
            shape.text = target


def delete_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def set_textbox(shape, text: str, font_size: int = 16, bold_first: bool = False) -> None:
    tf = shape.text_frame
    tf.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = line
        run.font.size = Pt(font_size)
        if bold_first and i == 0:
            run.font.bold = True


def rebuild_slide8(slide, figure_path: Path) -> None:
    # Remove duplicate cards / manual bars / old takeaways.
    removable = []
    for shape in slide.shapes:
        name = shape.name
        if name.startswith(("Rounded Rectangle", "TextBox 9", "TextBox 10", "TextBox 11",
                            "TextBox 12", "TextBox 13", "TextBox 16", "TextBox 17",
                            "TextBox 20", "TextBox 21", "TextBox 24")):
            removable.append(shape)
        elif name in {f"Rounded Rectangle {i}" for i in range(4, 28)}:
            removable.append(shape)
    # Safer: remove everything except title/header/footer shapes 1-3 and 28-30.
    keep = {"TextBox 1", "TextBox 2", "Rectangle 3", "Rectangle 28", "TextBox 29", "TextBox 30"}
    for shape in list(slide.shapes):
        if shape.name not in keep:
            delete_shape(shape)

    slide.shapes.add_picture(
        str(figure_path),
        Inches(0.55),
        Inches(1.35),
        width=Inches(7.35),
    )

    box = slide.shapes.add_textbox(Inches(8.05), Inches(1.45), Inches(4.35), Inches(4.55))
    tf = box.text_frame
    tf.word_wrap = True
    bullets = [
        "Key takeaways (991-question full test):",
        "14B zero-shot 83.55% → distilled mean 88.67% (+5.12 pp); best 89.10%",
        "7B improves 76.49% → 85.60% (+9.11 pp)",
        "Best 14B exceeds DeepSeek-V3 teacher 87.18% (+1.92 pp)",
        "Stable across 3 seeds: 88.40%–89.10% (range 0.70 pp)",
        "Figure 1 summarizes the main benchmark comparison.",
    ]
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0 if i == 0 else 1
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = line
        run.font.size = Pt(15 if i == 0 else 13)
        if i == 0:
            run.font.bold = True


def update_title_slide(slide) -> None:
    replace_text_in_slide(
        slide,
        "Presenter: Tianyuan Chen",
        "Presenter: Tianyuan Chen\nCo-authors: Andrew Chi-Chung Cheng, Harris Sik-Ho Tsang, Tony Yu-Lin Zhu, Xiaoxing Yang, Wai-Lun Lo",
    )
    replace_text_in_slide(
        slide,
        "Supervisor: Dr. Richard Tai-Chiu Hsung",
        "Supervisor: Dr. Richard Tai-Chiu Hsung*\nCollaborators: Billy Hon Wing Chiu (Lingnan); R. Chau, Z. Gong, K. Mao, W. Lam (HKU Dentistry)",
    )
    for shape in slide.shapes:
        if shape.has_text_frame and "Co-authors:" in shape.text:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
        if shape.has_text_frame and "Collaborators:" in shape.text:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(13)


def main() -> None:
    shutil.copy2(SRC, BACKUP)
    figure_path = extract_paper_figure()

    prs = Presentation(str(SRC))

    # Footer page numbers.
    for idx, slide in enumerate(prs.slides, start=1):
        if idx >= 2:
            set_footer(slide, idx)

    # Slide 5 formula.
    s5 = prs.slides[4]
    replace_text_in_slide(
        s5,
        "L = alpha * KL(p_T || p_S) + (1 - alpha) * CE",
        "L = α · KL(p_T || p_S) + (1 − α) · CE",
    )
    replace_text_in_slide(
        s5,
        "Main setting: one Stage 1 epoch with alpha = 0.35 for the best 14B run.",
        "Main setting: one Stage 1 epoch with α = 0.35 for the best 14B run.",
    )
    set_footer(s5, 5)

    # Slide 7 setup.
    s7 = prs.slides[6]
    replace_text_in_slide(s7, "6,591 Questions", "6,590 Questions")
    replace_text_in_slide(
        s7,
        "Stage 2 is tested but not required for the best 14B result",
        "Teacher–GT disagree on ~12.2% of training items\nStage 2 is tested but not required for the best 14B result",
    )

    # Slide 8 results.
    rebuild_slide8(prs.slides[7], figure_path)

    # Slide 10 discussion scope.
    s10 = prs.slides[9]
    replace_text_in_slide(
        s10,
        "Decision-space transfer can be a practical design rule",
        "Decision-space transfer can be a practical design rule\nScope: fixed-choice MCQ only; not for open-ended clinical dialogue",
    )

    # Slide 1 authors.
    update_title_slide(prs.slides[0])

    prs.save(str(OUT))
    print(f"Backup: {BACKUP}")
    print(f"Updated: {OUT}")


if __name__ == "__main__":
    main()
