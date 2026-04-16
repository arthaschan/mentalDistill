#!/usr/bin/env python3
"""Generate a single PPT slide: Plan A vs Plan B comparison."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Background ──
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ── Colours ──
DARK   = RGBColor(0x2D, 0x3A, 0x4A)
ACCENT = RGBColor(0x00, 0x70, 0xC0)
GREEN  = RGBColor(0x00, 0x80, 0x50)
RED    = RGBColor(0xC0, 0x00, 0x00)
GRAY   = RGBColor(0x66, 0x66, 0x66)
HDR_BG = RGBColor(0x00, 0x70, 0xC0)
HDR_FG = RGBColor(0xFF, 0xFF, 0xFF)
ROW_A  = RGBColor(0xF2, 0xF7, 0xFC)
ROW_B  = RGBColor(0xFF, 0xFF, 0xFF)

# ══════════════════════════════════════════════════
# Title
# ══════════════════════════════════════════════════
tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.55))
tf = tx.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Module 15 — 全量 CMExam 重训练对比 (Plan A → Plan B)"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = DARK

# subtitle
tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(12), Inches(0.35))
tf2 = tx2.text_frame
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Plan A: 仅牙科 672 样本训练 / 83 题测试 → Plan B: 全科 4608 样本训练 / 991 题(+125 牙科子集)测试"
r2.font.size = Pt(12); r2.font.color.rgb = GRAY

# ══════════════════════════════════════════════════
# Helper: add a table
# ══════════════════════════════════════════════════
def add_table(slide, left, top, width, height, rows, cols, data, title, col_widths=None):
    """data = list of list of (text, bold, color)"""
    # title
    tx = slide.shapes.add_textbox(left, top - Inches(0.35), width, Inches(0.35))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = ACCENT

    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

    for ri in range(rows):
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            text, bold, color = data[ri][ci]
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = text
            run.font.size = Pt(10)
            run.font.bold = bold
            run.font.color.rgb = color if color else DARK

            # header row fill
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HDR_BG
                run.font.color.rgb = HDR_FG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_A if ri % 2 == 1 else ROW_B

    return tbl_shape

# ══════════════════════════════════════════════════
# 7B Table
# ══════════════════════════════════════════════════
hdr = [("", True, None), ("Seed 11", True, None), ("Seed 42", True, None), ("Seed 7/7", True, None), ("均值 ± σ", True, None), ("提升", True, None)]
data_7b = [
    hdr,
    [("7B Zero-shot", True, GRAY),   ("—", False, GRAY), ("—", False, GRAY), ("—", False, GRAY), ("76.49% (全) / 68.80% (牙)", False, GRAY), ("—", False, GRAY)],
    [("Plan A (牙科672)", True, None),("81.93%", False, None), ("72.29%", False, None), ("77.11%", False, None), ("76.87% ± 3.09pp (83题)", False, None), ("—", False, GRAY)],
    [("Plan B 全科 (991题)", True, GREEN), ("86.28%", False, GREEN), ("85.57%", False, GREEN), ("84.96%", False, GREEN), ("85.60% ± 0.54pp", True, GREEN), ("+9.11pp", True, GREEN)],
    [("Plan B 牙科 (125题)", True, None), ("76.00%", False, None), ("73.60%", False, None), ("71.20%", False, None), ("73.60%", False, None), ("+4.80pp", False, GREEN)],
    [("教师 DeepSeek-V3", True, GRAY),("—", False, GRAY), ("—", False, GRAY), ("—", False, GRAY), ("87.18% (全) / 79.20% (牙)", False, GRAY), ("—", False, GRAY)],
]

add_table(slide,
    left=Inches(0.3), top=Inches(1.35),
    width=Inches(6.2), height=Inches(1.9),
    rows=6, cols=6, data=data_7b,
    title="Qwen2.5-7B 学生 (两阶段蒸馏)",
    col_widths=[Inches(1.3), Inches(0.8), Inches(0.8), Inches(0.8), Inches(1.7), Inches(0.8)])

# ══════════════════════════════════════════════════
# 14B Table
# ══════════════════════════════════════════════════
hdr14 = [("", True, None), ("Seed 11", True, None), ("Seed 42", True, None), ("Seed 8/8", True, None), ("均值 ± σ", True, None), ("提升", True, None)]
data_14b = [
    hdr14,
    [("14B Zero-shot", True, GRAY),   ("—", False, GRAY), ("—", False, GRAY), ("—", False, GRAY), ("83.25% (全) / 79.20% (牙)", False, GRAY), ("—", False, GRAY)],
    [("Plan A (牙科672)", True, None),("84.34%", False, None), ("83.13%", False, None), ("79.52%", False, None), ("82.33% ± 1.99pp (83题)", False, None), ("—", False, GRAY)],
    [("Plan B 全科 (991题)", True, GREEN), ("88.50%", False, GREEN), ("88.40%", False, GREEN), ("89.10%", False, GREEN), ("88.67% ± 0.31pp", True, GREEN), ("+5.42pp", True, GREEN)],
    [("Plan B 牙科 (125题)", True, None), ("79.20%", False, None), ("80.00%", False, None), ("78.40%", False, None), ("79.20%", False, None), ("+0.00pp", False, RED)],
    [("教师 DeepSeek-V3", True, GRAY),("—", False, GRAY), ("—", False, GRAY), ("—", False, GRAY), ("87.18% (全) / 79.20% (牙)", False, GRAY), ("—", False, GRAY)],
]

add_table(slide,
    left=Inches(6.8), top=Inches(1.35),
    width=Inches(6.2), height=Inches(1.9),
    rows=6, cols=6, data=data_14b,
    title="Qwen2.5-14B 学生 (仅 Stage-1 蒸馏)",
    col_widths=[Inches(1.3), Inches(0.8), Inches(0.8), Inches(0.8), Inches(1.7), Inches(0.8)])

# ══════════════════════════════════════════════════
# Conclusions box
# ══════════════════════════════════════════════════
box_left = Inches(0.3)
box_top  = Inches(3.65)
box_w    = Inches(12.7)
box_h    = Inches(3.55)

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, box_top, box_w, box_h)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
shape.line.color.rgb = ACCENT
shape.line.width = Pt(1.2)

tf = shape.text_frame
tf.word_wrap = True

conclusions = [
    ("核心发现", True, ACCENT, 14),
    ("", False, DARK, 6),
    ("1. 全量训练显著提升：7B 全科准确率 85.60%（+9.11pp），14B 全科 88.67%（+5.42pp），大数据量 × 多学科覆盖为蒸馏带来实质性增益。", False, DARK, 11),
    ("2. 14B 学生超越教师：14B 均值 88.67% > DeepSeek-V3 教师 87.18%，验证了「强学生 + 大数据蒸馏可超越教师」的假设。", False, DARK, 11),
    ("3. 种子方差大幅压缩：测试集从 83→991 题后，7B σ 从 3.09→0.54pp（-83%），14B σ 从 1.99→0.31pp（-84%），结论更可靠。", False, DARK, 11),
    ("4. 牙科子集提升有限：14B 牙科 79.20% 与 zero-shot 持平（已触及教师天花板 79.20%），7B 牙科 +4.80pp 仍有改善。", False, DARK, 11),
    ("", False, DARK, 6),
    ("实验设置", True, ACCENT, 14),
    ("", False, DARK, 6),
    ("• 数据集：CMExam 6591 道单选题（去重后），按难度分层 seed=2026 切分：Train 4608 / Val 991 / Test 991（含 125 牙科子集）", False, DARK, 10),
    ("• 教师：DeepSeek-V3 API（全科 87.18%，牙科 79.20%）", False, DARK, 10),
    ("• 7B 配置：两阶段（Stage1 KL α=0.35 + Stage2 GT-SFT），LR 1.2e-4，LoRA r=16/α=32，smooth_eps=0.25", False, DARK, 10),
    ("• 14B 配置：仅 Stage-1（Stage2 对强学生有害），LR 1e-4，LoRA r=16/α=32", False, DARK, 10),
    ("• 每组 3 seed 取均值 ± 标准差", False, DARK, 10),
]

for i, (text, bold, color, size) in enumerate(conclusions):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    p.space_after = Pt(1)

# ── Save ──
out = "/home/student/arthas/mentalDistill/docs/module15_comparison.pptx"
prs.save(out)
print(f"Saved → {out}")
