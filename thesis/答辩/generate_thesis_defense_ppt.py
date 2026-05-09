from __future__ import annotations

from io import BytesIO
from pathlib import Path
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PPT = ROOT / "RAI+601+毕业设计汇报 #10 陈天元256360231.pptx"
OUTPUT_PPT = ROOT / "答辩" / "论文汇报PPT_模板风格版.pptx"

SLIDE_WIDTH = 12179300
SLIDE_HEIGHT = 6858000

FONT_CN = "Noto Sans CJK SC"
FONT_EN = "Calibri"

RED = RGBColor(155, 31, 35)
RED_DARK = RGBColor(122, 20, 23)
TEXT = RGBColor(48, 48, 48)
TEXT_LIGHT = RGBColor(92, 92, 92)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(31, 96, 156)
TEAL = RGBColor(0, 138, 150)
GREEN = RGBColor(28, 145, 93)
ORANGE = RGBColor(220, 116, 26)
LIGHT_BLUE = RGBColor(231, 239, 248)
LIGHT_GRAY = RGBColor(247, 247, 247)


def set_text_style(paragraph, font_size, color=TEXT, bold=False, font_name=FONT_CN):
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=20,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name=FONT_CN,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    set_text_style(p, font_size, color=color, bold=bold, font_name=font_name)
    return box


def add_bullets(slide, left, top, width, height, items, *, font_size=22, color=TEXT_LIGHT, level0_indent=0.28):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10)
        p.space_before = Pt(0)
        p.left_margin = Inches(level0_indent)
        set_text_style(p, font_size, color=color)
    return box


def add_footer(slide, index, total, title):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, SLIDE_HEIGHT - Inches(0.5), SLIDE_WIDTH, Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()

    add_textbox(
        slide,
        Inches(4.2),
        SLIDE_HEIGHT - Inches(0.39),
        Inches(5.2),
        Inches(0.25),
        title,
        font_size=16,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        SLIDE_WIDTH - Inches(0.8),
        SLIDE_HEIGHT - Inches(0.39),
        Inches(0.55),
        Inches(0.25),
        f"{index}/{total}",
        font_size=16,
        color=WHITE,
        align=PP_ALIGN.RIGHT,
        font_name=FONT_EN,
    )


def add_title(slide, number, title, *, title_color=TEXT):
    add_textbox(slide, Inches(0.7), Inches(0.42), Inches(0.9), Inches(0.5), number, font_size=30, color=RED, font_name=FONT_EN)
    add_textbox(slide, Inches(1.75), Inches(0.37), Inches(6.8), Inches(0.6), title, font_size=30, color=title_color, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(1.18), Inches(2.05), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()


def add_round_box(slide, left, top, width, height, *, fill, line, radius_text=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    if radius_text:
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = radius_text
        set_text_style(p, 18, color=TEXT)
    return shape


def add_card(slide, left, top, width, height, title, body, *, line_color, title_color=None):
    card = add_round_box(slide, left, top, width, height, fill=WHITE, line=line_color)
    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.space_after = Pt(6)
    set_text_style(p1, 18, color=title_color or line_color, bold=True)

    for line in body:
        p = tf.add_paragraph()
        p.text = line
        p.bullet = False
        p.space_after = Pt(4)
        set_text_style(p, 14, color=TEXT_LIGHT)
    return card


def add_metric_card(slide, left, top, width, height, title, value, note, *, accent):
    card = add_round_box(slide, left, top, width, height, fill=LIGHT_GRAY, line=accent)
    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    set_text_style(p1, 15, color=accent, bold=True)
    p2 = tf.add_paragraph()
    p2.text = value
    p2.space_before = Pt(4)
    set_text_style(p2, 24, color=TEXT, bold=True, font_name=FONT_EN)
    p3 = tf.add_paragraph()
    p3.text = note
    p3.space_before = Pt(2)
    set_text_style(p3, 12, color=TEXT_LIGHT)
    return card


def add_logo(slide):
    if not TEMPLATE_PPT.exists():
        return
    with zipfile.ZipFile(TEMPLATE_PPT) as archive:
        data = archive.read("ppt/media/image1.png")
    slide.shapes.add_picture(BytesIO(data), Inches(0.35), Inches(0.15), width=Inches(2.3))


def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_logo(slide)
    add_textbox(slide, Inches(3.35), Inches(1.35), Inches(9.1), Inches(1.0), "基于知识蒸馏的牙科选择题自动答题系统", font_size=34, color=TEXT, bold=True)
    add_textbox(slide, Inches(4.9), Inches(2.45), Inches(6.6), Inches(0.4), "Knowledge Distillation for Dental MCQ Answering", font_size=20, color=RGBColor(132, 132, 132), font_name=FONT_EN)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(5.1), Inches(3.1), Inches(3.4), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    add_textbox(slide, Inches(5.45), Inches(3.55), Inches(5.7), Inches(0.4), "NVIDIA H100 NVL 95GB  |  CMExam 牙科医师考试选择题", font_size=18, color=TEXT_LIGHT)
    add_textbox(slide, Inches(5.15), Inches(4.0), Inches(5.8), Inches(0.45), "核心结果：89.10%（14B 学生） > 87.18%（DeepSeek-V3 教师）", font_size=20, color=GREEN, bold=True)

    info = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, SLIDE_HEIGHT - Inches(1.9), SLIDE_WIDTH, Inches(1.75))
    info.fill.solid()
    info.fill.fore_color.rgb = RED
    info.line.fill.background()
    add_textbox(slide, Inches(2.35), Inches(5.72), Inches(2.6), Inches(0.35), "专业：应用人工智能", font_size=22, color=WHITE, bold=True)
    add_textbox(slide, Inches(7.15), Inches(5.72), Inches(3.3), Inches(0.35), "导师：熊体超博士", font_size=22, color=WHITE, bold=True)
    add_textbox(slide, Inches(2.35), Inches(6.4), Inches(2.6), Inches(0.35), "学号：256360231", font_size=22, color=WHITE, bold=True, font_name=FONT_EN)
    add_textbox(slide, Inches(7.15), Inches(6.4), Inches(2.8), Inches(0.35), "答辩人：陈天元", font_size=22, color=WHITE, bold=True)

    bottom = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, SLIDE_HEIGHT - Inches(0.48), SLIDE_WIDTH, Inches(0.43))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = RED_DARK
    bottom.line.fill.background()
    add_textbox(slide, Inches(5.5), Inches(7.07), Inches(2.7), Inches(0.2), "2026 年 5 月", font_size=16, color=WHITE, align=PP_ALIGN.CENTER)


def build_outline(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, Inches(0.9), Inches(0.6), Inches(2.2), Inches(0.6), "汇报提纲", font_size=30, color=TEXT, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(1.18), Inches(1.55), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    items = [
        "1. 研究背景与问题定义",
        "2. 研究目标、数据与实验范围",
        "3. 核心方法：Choice-Head 两阶段蒸馏",
        "4. 实验体系与设置",
        "5. 核心结果：学生超过教师",
        "6. 关键发现与理论解释",
        "7. 应用价值、局限与未来工作",
        "8. 结论与答辩总结",
    ]
    add_bullets(slide, Inches(1.0), Inches(1.7), Inches(7.3), Inches(4.9), items, font_size=22, color=TEXT, level0_indent=0.18)
    add_footer(slide, 2, total, "基于知识蒸馏的牙科选择题自动答题系统")


def build_background(slide, total):
    add_title(slide, "01", "研究背景与问题定义", title_color=RED)
    box = add_round_box(slide, Inches(0.8), Inches(1.65), Inches(11.3), Inches(1.15), fill=LIGHT_BLUE, line=TEAL)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "核心问题：能否通过知识蒸馏，将大型教师模型的答题能力迁移到轻量级学生模型（7B~14B），在保持可部署性的同时最大化准确率？"
    set_text_style(p, 21, color=BLUE, bold=True)

    add_textbox(slide, Inches(0.9), Inches(3.2), Inches(2.0), Inches(0.5), "研究动机", font_size=24, color=TEXT, bold=True)
    bullets = [
        "牙科知识服务需要即时、低门槛、可解释的辅助能力",
        "商业级教师模型准确率高，但显存、算力和调用成本高",
        "传统全词表蒸馏对五选一任务存在明显冗余",
        "研究目标不是做更大模型，而是把高性能能力做成可部署系统",
    ]
    add_bullets(slide, Inches(0.95), Inches(3.75), Inches(10.5), Inches(2.6), bullets, font_size=21)
    add_footer(slide, 3, total, "基于知识蒸馏的牙科选择题自动答题系统")


def build_setup(slide, total):
    add_title(slide, "02", "研究目标、数据与实验范围")
    add_metric_card(slide, Inches(0.85), Inches(1.65), Inches(2.45), Inches(1.2), "数据范围", "6,591 题", "全量 CMExam 单选题，含牙科子集", accent=TEAL)
    add_metric_card(slide, Inches(3.5), Inches(1.65), Inches(2.45), Inches(1.2), "实验规模", "21 组", "覆盖 00–20 模块的系统性实验", accent=ORANGE)
    add_metric_card(slide, Inches(6.15), Inches(1.65), Inches(2.45), Inches(1.2), "学生模型", "7B / 14B", "Qwen2.5 系列 LoRA 微调", accent=GREEN)
    add_metric_card(slide, Inches(8.8), Inches(1.65), Inches(2.45), Inches(1.2), "部署约束", "45× 压缩", "从 671B 教师迁移到 14B 学生", accent=BLUE)

    left = add_round_box(slide, Inches(0.85), Inches(3.1), Inches(5.25), Inches(2.65), fill=WHITE, line=TEAL)
    tf_left = left.text_frame
    tf_left.clear()
    p = tf_left.paragraphs[0]
    p.text = "评测设置"
    set_text_style(p, 20, color=TEAL, bold=True)
    for line in [
        "小规模实验：672/74/83（train/val/test）牙科固定划分",
        "大规模验证：4608/991/991 全量重分割 + 125 题牙科子集",
        "教师：DeepSeek-V3、Doubao、Llama-3.3-70B、Qwen2.5-32B",
        "硬件：NVIDIA H100 NVL 95GB，BF16 混合精度训练",
    ]:
        para = tf_left.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)

    right = add_round_box(slide, Inches(6.35), Inches(3.1), Inches(5.0), Inches(2.65), fill=WHITE, line=ORANGE)
    tf_right = right.text_frame
    tf_right.clear()
    p = tf_right.paragraphs[0]
    p.text = "模型参数规模"
    set_text_style(p, 20, color=ORANGE, bold=True)
    for line in [
        "学生：Qwen2.5-7B / 14B，覆盖 7B 到 14B 容量段",
        "本地教师：Qwen2.5-32B（32B）、Llama-3.3-70B（70B）",
        "主要 API 教师：DeepSeek-V3（671B 总参数 / 37B 激活）",
        "Doubao、Kimi 为闭源模型，参数规模未公开",
    ]:
        para = tf_right.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)
    add_footer(slide, 4, total, "基于知识蒸馏的牙科选择题自动答题系统")


def build_method(slide, total):
    add_title(slide, "03", "核心方法：Choice-Head 两阶段蒸馏")
    add_card(slide, Inches(0.8), Inches(1.55), Inches(2.25), Inches(1.65), "传统全词表蒸馏", ["151,936 维 logits", "显存与计算代价高", "API 教师难以支持"], line_color=RGBColor(120, 120, 120), title_color=TEXT)
    add_card(slide, Inches(3.25), Inches(1.55), Inches(2.35), Inches(1.65), "Choice-Head 蒸馏", ["只蒸馏 A/B/C/D/E", "监督维度压缩到 5 维", "天然兼容黑盒教师"], line_color=TEAL)
    add_card(slide, Inches(5.8), Inches(1.55), Inches(2.35), Inches(1.65), "Stage 1", ["Choice-Head KL + CE", "1 epoch，α=0.35", "先学教师偏好结构"], line_color=ORANGE)
    add_card(slide, Inches(8.35), Inches(1.55), Inches(3.0), Inches(1.65), "Stage 2", ["GT SFT 精校", "2–5 epochs", "用标准答案纠正教师误差"], line_color=BLUE)

    add_textbox(slide, Inches(0.95), Inches(3.55), Inches(4.8), Inches(0.5), "关键设计思想", font_size=24, color=RED, bold=True)
    left = add_round_box(slide, Inches(0.8), Inches(4.05), Inches(5.5), Inches(2.2), fill=LIGHT_GRAY, line=TEAL)
    tf_left = left.text_frame
    tf_left.clear()
    p = tf_left.paragraphs[0]
    p.text = "Stage 1：Choice-Head KL 蒸馏"
    set_text_style(p, 20, color=TEAL, bold=True)
    for line in [
        "仅提取 A/B/C/D/E 五个选项上的概率分布",
        "损失：α·KL(pT∥pS) + (1-α)·CE",
        "典型显存约 22GB，可使用 API-only 教师",
    ]:
        para = tf_left.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)

    right = add_round_box(slide, Inches(6.65), Inches(4.05), Inches(4.7), Inches(2.2), fill=LIGHT_GRAY, line=BLUE)
    tf_right = right.text_frame
    tf_right.clear()
    p = tf_right.paragraphs[0]
    p.text = "Stage 2：GT SFT 精校"
    set_text_style(p, 20, color=BLUE, bold=True)
    for line in [
        "继承 Stage 1 的 LoRA 权重",
        "用标准答案对教师错误进行校准",
        "对弱学生更稳定，对 14B 强学生不一定持续增益",
    ]:
        para = tf_right.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)
    add_footer(slide, 5, total, "基于知识蒸馏的牙科选择题自动答题系统")


def build_experiment_matrix(slide, total):
    add_title(slide, "04", "实验体系与设置")
    add_card(slide, Inches(0.8), Inches(1.55), Inches(2.15), Inches(1.55), "A. 基线建立", ["GT SFT 基线", "7B：77.11%", "明确学生起点"], line_color=RGBColor(90, 90, 90), title_color=TEXT)
    add_card(slide, Inches(3.05), Inches(1.55), Inches(2.35), Inches(1.55), "B. 单教师蒸馏", ["白盒 Logit-KL", "黑盒 Choice-Head", "比较教师形态与参数规模"], line_color=TEAL)
    add_card(slide, Inches(5.55), Inches(1.55), Inches(2.35), Inches(1.55), "C. 多教师优化", ["静态融合", "一致性过滤", "多数票集成"], line_color=ORANGE)
    add_card(slide, Inches(8.05), Inches(1.55), Inches(1.55), Inches(1.55), "D. 范式拓展", ["CoT 蒸馏", "α-散度", "边界过滤"], line_color=RED)
    add_card(slide, Inches(9.8), Inches(1.55), Inches(1.55), Inches(1.55), "E. 容量升级", ["14B 学生", "跨架构蒸馏", "全量重分割"], line_color=GREEN)

    left = add_round_box(slide, Inches(0.8), Inches(3.55), Inches(5.25), Inches(2.35), fill=WHITE, line=TEAL)
    tf_left = left.text_frame
    tf_left.clear()
    p = tf_left.paragraphs[0]
    p.text = "实验设计特点"
    set_text_style(p, 20, color=TEAL, bold=True)
    for line in [
        "从单教师到多教师、从 7B 到 14B、从 83 题到 991 题逐层推进",
        "包含正例、负例与失败实验，不只保留最优结果",
        "所有模块保留统一运行脚本，强调可复现性",
    ]:
        para = tf_left.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)

    right = add_round_box(slide, Inches(6.35), Inches(3.55), Inches(5.0), Inches(2.35), fill=WHITE, line=ORANGE)
    tf_right = right.text_frame
    tf_right.clear()
    p = tf_right.paragraphs[0]
    p.text = "关键可比性控制"
    set_text_style(p, 20, color=ORANGE, bold=True)
    for line in [
        "模块之间尽量复用同一学生架构与训练超参",
        "通过多 seed 与大测试集缓解 83 题小样本波动",
        "把理论分析放在经验结论之后，避免先验解释替代实证比较",
    ]:
        para = tf_right.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)
    add_footer(slide, 6, total, "基于知识蒸馏的牙科选择题自动答题系统")


def build_main_results(slide, total):
    add_title(slide, "05", "核心结果：学生超过教师")
    add_metric_card(slide, Inches(0.85), Inches(1.55), Inches(2.45), Inches(1.3), "DeepSeek-V3 教师", "87.18%", "991 题全量测试集准确率", accent=RED)
    add_metric_card(slide, Inches(3.55), Inches(1.55), Inches(2.45), Inches(1.3), "14B 学生最佳", "89.10%", "Module 15，seed 8", accent=GREEN)
    add_metric_card(slide, Inches(6.25), Inches(1.55), Inches(2.45), Inches(1.3), "14B 三种子均值", "88.67%", "同样超过教师表现", accent=BLUE)
    add_metric_card(slide, Inches(8.95), Inches(1.55), Inches(2.45), Inches(1.3), "参数压缩比", "45×", "671B 教师到 14B 学生", accent=ORANGE)

    left = add_round_box(slide, Inches(0.85), Inches(3.2), Inches(5.3), Inches(2.45), fill=WHITE, line=GREEN)
    tf_left = left.text_frame
    tf_left.clear()
    p = tf_left.paragraphs[0]
    p.text = "最重要的结论"
    set_text_style(p, 22, color=GREEN, bold=True)
    for line in [
        "学生模型不只是接近教师，而是在任务结构更对齐的训练目标下可以超过教师",
        "对 14B 强学生而言，最有效的并不一定是更复杂流程，而是更精准的蒸馏目标",
        "大测试集验证后，结论在统计上更稳健，不依赖少量题目的偶然波动",
    ]:
        para = tf_left.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)

    right = add_round_box(slide, Inches(6.45), Inches(3.2), Inches(4.9), Inches(2.45), fill=LIGHT_GRAY, line=RED)
    tf_right = right.text_frame
    tf_right.clear()
    p = tf_right.paragraphs[0]
    p.text = "代表性结果串联"
    set_text_style(p, 20, color=RED, bold=True)
    for line in [
        "7B GT SFT 基线：77.11%",
        "7B Choice-Head 蒸馏最佳：81.93%",
        "14B DeepSeek 蒸馏最佳：84.34%（83 题牙科）",
        "14B 全量重分割最佳：89.10%（991 题）",
    ]:
        para = tf_right.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)
    add_footer(slide, 7, total, "基于知识蒸馏的牙科选择题自动答题系统")


def build_result_details(slide, total):
    add_title(slide, "06", "结果细化：什么条件下蒸馏有效？")
    add_card(slide, Inches(0.8), Inches(1.55), Inches(3.35), Inches(1.9), "白盒 vs 黑盒", ["白盒 Logit-KL：80.72%", "黑盒 Choice-Head：81.93%", "任务对齐优于全词表监督"], line_color=TEAL)
    add_card(slide, Inches(4.35), Inches(1.55), Inches(3.35), Inches(1.9), "Stage 2 的边界", ["7B 学生通常受益", "14B 学生均值反而下降约 1.6pp", "说明强学生会被过校准"], line_color=ORANGE)
    add_card(slide, Inches(7.9), Inches(1.55), Inches(3.35), Inches(1.9), "跨架构可行性", ["72.8% 的 Llama-70B 教师", "仍能蒸馏出 87.25% 的 14B 学生", "监督价值不只由教师准确率决定"], line_color=GREEN)

    bottom = add_round_box(slide, Inches(0.85), Inches(3.95), Inches(10.5), Inches(1.8), fill=LIGHT_BLUE, line=BLUE)
    tf = bottom.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "结论不是“蒸馏越复杂越好”，而是“监督信号是否与任务决策结构匹配”。"
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, 23, color=BLUE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "因此，本文的主要价值不只是一组更高分数，而是给出了一套可解释的设计原则：聚焦选项空间、保留教师偏好、避免无效监督维度。"
    p2.alignment = PP_ALIGN.CENTER
    set_text_style(p2, 16, color=TEXT_LIGHT)
    add_footer(slide, 8, total, "基于知识蒸馏的牙科选择题自动答题系统")


def build_inverted_u(slide, total):
    add_title(slide, "07", "关键发现一：教师不是越强越好")
    add_card(slide, Inches(0.9), Inches(1.7), Inches(3.1), Inches(2.7), "弱教师", ["代表：Kimi，61.45%", "错误过多，噪声大", "学生学习到的偏好不稳定"], line_color=RED)
    add_card(slide, Inches(4.15), Inches(1.7), Inches(4.0), Inches(2.7), "最佳区间", ["教师与 GT 不一致率约 5%–15%", "既保留结构信息，又不过度偏离真值", "是最适合迁移的蒸馏信号"], line_color=GREEN)
    add_card(slide, Inches(8.35), Inches(1.7), Inches(3.0), Inches(2.7), "过强教师", ["接近 one-hot", "暗知识不足", "难以提供有梯度价值的混淆结构"], line_color=ORANGE)

    box = add_round_box(slide, Inches(0.95), Inches(4.75), Inches(10.35), Inches(1.2), fill=WHITE, line=TEAL)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "教师质量与蒸馏收益呈倒 U 型关系。真正关键的不是教师“绝对有多强”，而是它是否保留了可迁移的结构性不确定性。"
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, 20, color=TEAL, bold=True)
    add_footer(slide, 9, total, "基于知识蒸馏的牙科选择题自动答题系统")


def build_information_geometry(slide, total):
    add_title(slide, "08", "关键发现二：真实 logprobs 更有价值")
    left = add_round_box(slide, Inches(0.8), Inches(1.55), Inches(5.25), Inches(3.0), fill=WHITE, line=BLUE)
    tf_left = left.text_frame
    tf_left.clear()
    p = tf_left.paragraphs[0]
    p.text = "为什么人工平滑标签不够？"
    set_text_style(p, 20, color=BLUE, bold=True)
    for line in [
        "人工平滑只在正确答案周围分配固定噪声，信息结构单一",
        "无法反映“最容易混淆的错误项”与“边界距离”",
        "因此在任务层面只提供了有限的软监督增益",
    ]:
        para = tf_left.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)

    right = add_round_box(slide, Inches(6.2), Inches(1.55), Inches(5.1), Inches(3.0), fill=WHITE, line=GREEN)
    tf_right = right.text_frame
    tf_right.clear()
    p = tf_right.paragraphs[0]
    p.text = "信息几何分析给出的证据"
    set_text_style(p, 20, color=GREEN, bold=True)
    for line in [
        "Fisher-Rao 与 α-散度分析显示：真实 logprobs 更接近连续概率流形",
        "真实标签的体积密度约为人工平滑标签的 2500 倍",
        "这解释了为何真实 logprobs 教师往往蒸馏效果更稳定",
    ]:
        para = tf_right.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)

    highlight = add_round_box(slide, Inches(1.4), Inches(4.95), Inches(9.3), Inches(0.9), fill=LIGHT_BLUE, line=TEAL)
    tf = highlight.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "结论：决定蒸馏效果的不只是“答对没有”，而是教师分布中是否保留了高价值的连续结构信息。"
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, 18, color=TEAL, bold=True)
    add_footer(slide, 10, total, "基于知识蒸馏的牙科选择题自动答题系统")


def build_application(slide, total):
    add_title(slide, "09", "应用价值与系统原型")
    left = add_round_box(slide, Inches(0.8), Inches(1.6), Inches(5.25), Inches(3.9), fill=WHITE, line=TEAL)
    tf_left = left.text_frame
    tf_left.clear()
    p = tf_left.paragraphs[0]
    p.text = "从论文到系统"
    set_text_style(p, 20, color=TEAL, bold=True)
    for line in [
        "仓库中已实现统一训练、评估、Web 推理与 Quiz 交互入口",
        "7B/14B 学生模型可以在普通 GPU 条件下完成部署演示",
        "适合患者教育、教学演示、标准化考试训练等轻量场景",
        "方法本身也可迁移到其他标准化多选知识评测任务",
    ]:
        para = tf_left.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)

    right = add_round_box(slide, Inches(6.3), Inches(1.6), Inches(5.0), Inches(3.9), fill=LIGHT_GRAY, line=ORANGE)
    tf_right = right.text_frame
    tf_right.clear()
    p = tf_right.paragraphs[0]
    p.text = "应用意义"
    set_text_style(p, 20, color=ORANGE, bold=True)
    for line in [
        "不是把大模型原样搬到端侧，而是围绕任务结构重构蒸馏目标",
        "把“能答题”转化为“能低成本稳定交付”",
        "为医疗大模型轻量化提供了可复现、可解释的实验范式",
    ]:
        para = tf_right.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)

    add_footer(slide, 11, total, "基于知识蒸馏的牙科选择题自动答题系统")


def build_future(slide, total):
    add_title(slide, "10", "局限性与未来工作")
    left = add_round_box(slide, Inches(0.85), Inches(1.65), Inches(5.15), Inches(3.95), fill=WHITE, line=RED)
    tf_left = left.text_frame
    tf_left.clear()
    p = tf_left.paragraphs[0]
    p.text = "当前局限"
    set_text_style(p, 21, color=RED, bold=True)
    for line in [
        "任务仍以标准化选择题为主，尚未覆盖开放问答与长文本推理",
        "教师谱系还不够密集，倒 U 型规律仍需更细粒度验证",
        "评估指标以准确率为主，缺少人工质量评审与临床专家反馈",
        "系统目前更多面向研究验证，而非直接临床部署",
    ]:
        para = tf_left.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)

    right = add_round_box(slide, Inches(6.2), Inches(1.65), Inches(5.15), Inches(3.95), fill=WHITE, line=GREEN)
    tf_right = right.text_frame
    tf_right.clear()
    p = tf_right.paragraphs[0]
    p.text = "未来工作"
    set_text_style(p, 21, color=GREEN, bold=True)
    for line in [
        "扩大数据与多种子验证，进一步提高结论的统计置信度",
        "研究教师路由、自适应 α-散度与动态样本选择策略",
        "从选择题扩展到开放医疗问答、多模态与真实交互任务",
        "继续推进端侧部署与真实用户反馈闭环",
    ]:
        para = tf_right.add_paragraph()
        para.text = line
        para.bullet = True
        para.left_margin = Inches(0.25)
        set_text_style(para, 16, color=TEXT_LIGHT)
    add_footer(slide, 12, total, "基于知识蒸馏的牙科选择题自动答题系统")


def build_ending(slide, total):
    add_textbox(slide, Inches(3.0), Inches(1.65), Inches(7.0), Inches(0.9), "结论与答辩总结", font_size=34, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(4.7), Inches(2.45), Inches(2.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    box = add_round_box(slide, Inches(1.2), Inches(3.0), Inches(10.0), Inches(2.2), fill=LIGHT_BLUE, line=TEAL)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "本文证明：当蒸馏目标与任务决策结构对齐时，小模型不仅能以更低成本部署，还有机会获得超过教师的任务表现。"
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, 22, color=BLUE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "感谢各位老师聆听，恳请批评指正。"
    p2.alignment = PP_ALIGN.CENTER
    set_text_style(p2, 24, color=GREEN, bold=True)
    add_footer(slide, 13, total, "基于知识蒸馏的牙科选择题自动答题系统")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    total = 13
    build_cover(prs)
    build_outline(prs, total)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_background(slide, total)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_setup(slide, total)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_method(slide, total)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_experiment_matrix(slide, total)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_main_results(slide, total)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_result_details(slide, total)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_inverted_u(slide, total)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_information_geometry(slide, total)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_application(slide, total)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_future(slide, total)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_ending(slide, total)

    OUTPUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPT)
    print(f"Generated: {OUTPUT_PPT}")


if __name__ == "__main__":
    main()